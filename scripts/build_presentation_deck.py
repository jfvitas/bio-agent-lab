from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "docs" / "presentations"
ASSET_DIR = OUTPUT_DIR / "assets"
SCREENSHOT_PATH = ASSET_DIR / "proteosphere_winui_home.png"
WIDE_LOGO_PATH = ROOT / "apps" / "PbdataWinUI" / "Assets" / "Wide310x150Logo.scale-200.png"
ICON_PATH = ROOT / "apps" / "PbdataWinUI" / "Assets" / "Square150x150Logo.scale-200.png"

BG = RGBColor(247, 249, 252)
NAVY = RGBColor(18, 33, 54)
TEAL = RGBColor(18, 191, 198)
BLUE = RGBColor(43, 106, 216)
GOLD = RGBColor(238, 184, 60)
CORAL = RGBColor(244, 116, 138)
SLATE = RGBColor(96, 107, 120)
LIGHT = RGBColor(230, 237, 244)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(36, 42, 49)
MINT = RGBColor(228, 247, 245)


def read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def load_metrics() -> dict:
    bootstrap = read_json(ROOT / "metadata" / "bootstrap_catalog" / "bootstrap_store_manifest.json")
    lifecycle = read_json(ROOT / "data" / "reports" / "source_lifecycle_report.json")
    training = read_json(ROOT / "custom_training_summary.json")
    splits = read_json(ROOT / "data" / "splits" / "metadata.json")
    audit = read_json(ROOT / "data" / "reports" / "screening_field_audit.json")
    return {
        "bootstrap_records": bootstrap.get("record_count", 0),
        "tracked_sources": lifecycle.get("summary", {}).get("tracked_sources", 0),
        "ready_sources": lifecycle.get("summary", {}).get("ready_sources", 0),
        "candidate_pool": training.get("candidate_pool_count", 0),
        "selected_examples": training.get("selected_count", 0),
        "selected_clusters": training.get("selected_receptor_clusters", 0),
        "selected_families": training.get("selected_metadata_families", 0),
        "selected_pathways": training.get("selected_pathway_groups", 0),
        "selected_folds": training.get("selected_fold_groups", 0),
        "mean_quality": training.get("mean_quality_score", 0.0),
        "split_strategy": splits.get("strategy", "unknown"),
        "split_sizes": splits.get("sizes", {}),
        "unsafe_fields": audit.get("unsafe_policy_field_count", 0),
        "restricted_fields": audit.get("restricted_policy_field_count", 0),
    }


def add_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, left, top, width, height, text, *, size=20, color=NAVY, bold=False,
                font_name="Aptos", align=PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bullets(slide, left, top, width, height, items, *, size=18, color=DARK) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = f"\u2022 {item}"
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)


def add_card(slide, left, top, width, height, *, fill_color=WHITE, line_color=LIGHT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    return shape


def add_stat_card(slide, left, top, width, height, title, value, body, accent):
    add_card(slide, left, top, width, height)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, left + Inches(0.16), top + Inches(0.18), width - Inches(0.32), Inches(0.24),
                title.upper(), size=10, color=accent, bold=True)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.44), width - Inches(0.32), Inches(0.45),
                value, size=22, color=NAVY, bold=True)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.92), width - Inches(0.32), height - Inches(1.02),
                body, size=11, color=SLATE)


def add_title_slide(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.26))
    header.fill.solid()
    header.fill.fore_color.rgb = TEAL
    header.line.fill.background()

    if WIDE_LOGO_PATH.exists():
        slide.shapes.add_picture(str(WIDE_LOGO_PATH), Inches(0.7), Inches(0.58), width=Inches(2.7))
    if ICON_PATH.exists():
        slide.shapes.add_picture(str(ICON_PATH), Inches(11.0), Inches(0.52), width=Inches(1.05))

    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(4.0), Inches(0.4),
                "PROTEOSPHERE", size=18, color=TEAL, bold=True)
    add_textbox(slide, Inches(0.72), Inches(1.95), Inches(11.0), Inches(0.75),
                "A protein intelligence platform for building trustworthy biological ML datasets and workflows", size=28, color=NAVY, bold=True)
    add_textbox(
        slide,
        Inches(0.72),
        Inches(2.95),
        Inches(8.8),
        Inches(1.35),
        "ProteoSphere is designed to overcome the limits of small, fragile, siloed protein-ML studies by unifying source packaging, curation, leakage-resistant dataset design, graph-aware modeling, and an intuitive desktop workflow.",
        size=19,
        color=SLATE,
    )

    add_stat_card(slide, Inches(0.72), Inches(5.02), Inches(2.15), Inches(1.5),
                  "Ready sources", f"{metrics['ready_sources']}/{metrics['tracked_sources']}",
                  "Source packages currently staged for local-first operation.", TEAL)
    add_stat_card(slide, Inches(3.03), Inches(5.02), Inches(2.15), Inches(1.5),
                  "Bootstrap index", f"{metrics['bootstrap_records']:,}",
                  "PDB records in the fast local planning layer.", BLUE)
    add_stat_card(slide, Inches(5.34), Inches(5.02), Inches(2.15), Inches(1.5),
                  "Candidate pool", f"{metrics['candidate_pool']:,}",
                  "Model-ready pairs available for selection.", GOLD)
    add_stat_card(slide, Inches(7.65), Inches(5.02), Inches(2.15), Inches(1.5),
                  "Selected set", f"{metrics['selected_examples']:,}",
                  "Current curated examples in the active training set.", CORAL)


def add_problem_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.7), Inches(0.52), Inches(11), Inches(0.45),
                "Why a different approach is needed", size=26, color=NAVY, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.0), Inches(10.8), Inches(0.58),
                "Most biological ML efforts still operate at a scale and workflow maturity that make their conclusions hard to generalize, compare, or reuse.", size=18, color=SLATE)

    add_card(slide, Inches(0.85), Inches(1.9), Inches(5.1), Inches(4.7))
    add_textbox(slide, Inches(1.1), Inches(2.15), Inches(4.3), Inches(0.3),
                "Limitations of typical small-scale studies", size=18, color=CORAL, bold=True)
    add_bullets(slide, Inches(1.1), Inches(2.55), Inches(4.2), Inches(3.5), [
        "Narrow, manually assembled datasets with limited provenance.",
        "Random or shallow splits that leak sequence, family, or assay similarity.",
        "Single-source assumptions that hide conflicts and coverage gaps.",
        "One-off graph or feature designs that do not scale beyond the original experiment.",
        "Tooling that is difficult to inspect, explain, or share with collaborators."
    ], size=16)

    add_card(slide, Inches(6.2), Inches(1.9), Inches(5.1), Inches(4.7), fill_color=MINT, line_color=TEAL)
    add_textbox(slide, Inches(6.45), Inches(2.15), Inches(4.3), Inches(0.3),
                "ProteoSphere response", size=18, color=TEAL, bold=True)
    add_bullets(slide, Inches(6.45), Inches(2.55), Inches(4.2), Inches(3.5), [
        "Local-first source packaging with explicit refresh and retention rules.",
        "Canonical planning, identity mapping, and decision-grade field auditing.",
        "Representative subset design plus leakage-resistant split governance.",
        "Multiple graph scopes and model pathways instead of a single narrow architecture.",
        "A GUI that makes the curation and modeling process understandable on screen."
    ], size=16)


def add_pipeline_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.7), Inches(0.52), Inches(11), Inches(0.45),
                "Pipeline strategy", size=26, color=NAVY, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.0), Inches(10.8), Inches(0.58),
                "The pipeline is intentionally layered so each stage adds scientific value, provenance, and reusability rather than just producing another flat table.", size=18, color=SLATE)

    stages = [
        ("1. Source packaging", "Stage broad local snapshots from structural, affinity, annotation, and pathway resources.", TEAL),
        ("2. Bootstrap planning", "Index PDB-level facts for fast selection, coverage review, and targeted refresh planning.", BLUE),
        ("3. Extraction + normalization", "Parse structures, assays, interfaces, and metadata into a coherent local schema.", GOLD),
        ("4. Screening + field audit", "Populate and validate decision-grade fields before they can influence curation policy.", CORAL),
        ("5. Training-set design", "Select representative examples while controlling cluster dominance and source bias.", TEAL),
        ("6. Split governance", "Build pair-aware, leakage-resistant train/validation/test partitions.", BLUE),
        ("7. Graph + feature packaging", "Export multiple structural graph designs plus engineered descriptor context.", GOLD),
        ("8. Model review + inference", "Train and compare suitable model paths, then surface explainable predictions.", CORAL),
    ]

    top = Inches(1.85)
    for index, (title, body, color) in enumerate(stages):
        row = index // 2
        col = index % 2
        left = Inches(0.9 + col * 6.1)
        y = top + Inches(row * 1.15)
        add_card(slide, left, y, Inches(5.25), Inches(0.95))
        marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.18), y + Inches(0.2), Inches(0.22), Inches(0.22))
        marker.fill.solid()
        marker.fill.fore_color.rgb = color
        marker.line.fill.background()
        add_textbox(slide, left + Inches(0.5), y + Inches(0.12), Inches(4.4), Inches(0.2), title, size=15, color=NAVY, bold=True)
        add_textbox(slide, left + Inches(0.5), y + Inches(0.38), Inches(4.55), Inches(0.35), body, size=12, color=SLATE)


def add_graph_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.7), Inches(0.52), Inches(11), Inches(0.45),
                "Graph and feature design principles", size=26, color=NAVY, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.0), Inches(10.8), Inches(0.58),
                "ProteoSphere is designed to support multiple structural representations rather than locking the project into one convenient graph story.", size=18, color=SLATE)

    add_card(slide, Inches(0.9), Inches(1.95), Inches(5.2), Inches(4.5))
    add_textbox(slide, Inches(1.15), Inches(2.18), Inches(4.3), Inches(0.3),
                "Supported graph design space", size=18, color=TEAL, bold=True)
    add_bullets(slide, Inches(1.15), Inches(2.58), Inches(4.2), Inches(3.2), [
        "Whole-protein graphs for global structural context.",
        "Interface-only graphs for contact-driven tasks.",
        "Shell and neighborhood graphs for focused local environments.",
        "Residue-level and atom-level graph constructions.",
        "Protein-ligand, protein-protein, and multimodal packaging."
    ], size=16)

    add_card(slide, Inches(6.35), Inches(1.95), Inches(5.2), Inches(4.5), fill_color=MINT, line_color=TEAL)
    add_textbox(slide, Inches(6.6), Inches(2.18), Inches(4.3), Inches(0.3),
                "Why this matters", size=18, color=BLUE, bold=True)
    add_bullets(slide, Inches(6.6), Inches(2.58), Inches(4.1), Inches(3.2), [
        "Different biological questions need different structural neighborhoods.",
        "Graph choice should be a controlled design variable, not an accidental default.",
        "Feature and graph coverage must stay tied to the actual dataset rows.",
        "Packaging should support both scalable baselines and richer graph-native models."
    ], size=16)


def add_dataset_slide(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.7), Inches(0.52), Inches(11), Inches(0.45),
                "Robust training-set creation", size=26, color=NAVY, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.0), Inches(10.8), Inches(0.58),
                "The goal is not to maximize row count at any cost. The goal is to build a set that is representative, diverse, auditable, and realistic for downstream evaluation.", size=18, color=SLATE)

    add_stat_card(slide, Inches(0.82), Inches(1.9), Inches(2.15), Inches(1.5),
                  "Candidate pool", f"{metrics['candidate_pool']:,}", "Model-ready pairs available for selection.", TEAL)
    add_stat_card(slide, Inches(3.12), Inches(1.9), Inches(2.15), Inches(1.5),
                  "Selected set", f"{metrics['selected_examples']:,}", "Curated examples in the active training subset.", BLUE)
    add_stat_card(slide, Inches(5.42), Inches(1.9), Inches(2.15), Inches(1.5),
                  "Clusters", f"{metrics['selected_clusters']:,}", "Receptor clusters represented in the selected set.", GOLD)
    add_stat_card(slide, Inches(7.72), Inches(1.9), Inches(2.15), Inches(1.5),
                  "Quality", f"{metrics['mean_quality']:.3f}", "Mean quality score for selected examples.", CORAL)

    add_card(slide, Inches(0.88), Inches(3.8), Inches(5.1), Inches(2.45))
    add_textbox(slide, Inches(1.12), Inches(4.05), Inches(4.3), Inches(0.3),
                "Selection principles", size=18, color=TEAL, bold=True)
    add_bullets(slide, Inches(1.12), Inches(4.45), Inches(4.2), Inches(1.55), [
        "Prefer representative coverage over near-duplicate clusters.",
        "Use explicit quality and provenance signals.",
        "Keep field-population quality visible before policy decisions are made."
    ], size=15)

    add_card(slide, Inches(6.22), Inches(3.8), Inches(5.1), Inches(2.45), fill_color=MINT, line_color=TEAL)
    add_textbox(slide, Inches(6.46), Inches(4.05), Inches(4.2), Inches(0.3),
                "Current diversity indicators", size=18, color=BLUE, bold=True)
    add_bullets(slide, Inches(6.46), Inches(4.45), Inches(4.1), Inches(1.55), [
        f"{metrics['selected_families']:,} metadata families represented.",
        f"{metrics['selected_pathways']:,} pathway groups represented.",
        f"{metrics['selected_folds']:,} fold groups represented."
    ], size=15)


def add_split_slide(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.7), Inches(0.52), Inches(11), Inches(0.45),
                "Leakage-resistant split strategy", size=26, color=NAVY, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.0), Inches(10.8), Inches(0.58),
                "Train/test splits are treated as a scientific design problem, not a final randomization step.", size=18, color=SLATE)

    train = metrics["split_sizes"].get("train", 0)
    val = metrics["split_sizes"].get("val", 0)
    test = metrics["split_sizes"].get("test", 0)

    add_card(slide, Inches(0.9), Inches(1.95), Inches(4.0), Inches(4.55))
    add_textbox(slide, Inches(1.15), Inches(2.18), Inches(3.3), Inches(0.3),
                "Split controls", size=18, color=TEAL, bold=True)
    add_bullets(slide, Inches(1.15), Inches(2.58), Inches(3.2), Inches(3.2), [
        "Pair-aware grouping to keep related pairs together.",
        "Sequence and family-aware grouping instead of naive row splitting.",
        "Mutation-cluster and source effects treated as leakage risks.",
        "Release split assignments kept explicit in the exported tables."
    ], size=16)

    add_card(slide, Inches(5.15), Inches(1.95), Inches(2.0), Inches(1.4))
    add_textbox(slide, Inches(5.4), Inches(2.15), Inches(1.4), Inches(0.2), "Train", size=16, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.4), Inches(2.45), Inches(1.4), Inches(0.3), f"{train:,}", size=28, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    add_card(slide, Inches(7.35), Inches(1.95), Inches(2.0), Inches(1.4))
    add_textbox(slide, Inches(7.6), Inches(2.15), Inches(1.4), Inches(0.2), "Validation", size=16, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.6), Inches(2.45), Inches(1.4), Inches(0.3), f"{val:,}", size=28, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    add_card(slide, Inches(9.55), Inches(1.95), Inches(2.0), Inches(1.4))
    add_textbox(slide, Inches(9.8), Inches(2.15), Inches(1.4), Inches(0.2), "Test", size=16, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.8), Inches(2.45), Inches(1.4), Inches(0.3), f"{test:,}", size=28, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    add_card(slide, Inches(5.15), Inches(3.7), Inches(6.4), Inches(2.8), fill_color=MINT, line_color=TEAL)
    add_textbox(slide, Inches(5.4), Inches(3.95), Inches(5.8), Inches(0.3),
                "What makes this stronger than typical practice", size=18, color=BLUE, bold=True)
    add_bullets(slide, Inches(5.4), Inches(4.35), Inches(5.7), Inches(1.9), [
        f"Current strategy: {metrics['split_strategy']}.",
        "The split logic is built to defend against inflated benchmark scores.",
        "The platform keeps the split story visible in both CLI and GUI outputs."
    ], size=16)


def add_model_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.7), Inches(0.52), Inches(11), Inches(0.45),
                "Model architecture strategy", size=26, color=NAVY, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.0), Inches(10.8), Inches(0.58),
                "ProteoSphere is not built around a single favorite model. It is built to compare and justify model paths based on the evidence available in the current dataset.", size=18, color=SLATE)

    add_card(slide, Inches(0.9), Inches(1.95), Inches(5.15), Inches(4.55))
    add_textbox(slide, Inches(1.15), Inches(2.18), Inches(4.3), Inches(0.3),
                "Supported modeling philosophy", size=18, color=TEAL, bold=True)
    add_bullets(slide, Inches(1.15), Inches(2.58), Inches(4.2), Inches(3.15), [
        "Use strong tabular baselines where graph coverage is weak.",
        "Escalate to graph-native or hybrid models when graph coverage justifies it.",
        "Treat recommendation, runtime readiness, and actual training behavior as one consistent system.",
        "Preserve provenance for recommended family, fallback behavior, and saved runs."
    ], size=16)

    add_card(slide, Inches(6.25), Inches(1.95), Inches(5.15), Inches(4.55), fill_color=MINT, line_color=TEAL)
    add_textbox(slide, Inches(6.5), Inches(2.18), Inches(4.3), Inches(0.3),
                "Benefits over one-off model studies", size=18, color=BLUE, bold=True)
    add_bullets(slide, Inches(6.5), Inches(2.58), Inches(4.15), Inches(3.15), [
        "Architecture choice is grounded in dataset readiness, not wishful complexity.",
        "Graph, feature, and runtime constraints are surfaced before training starts.",
        "Saved outputs are easier to compare, reload, explain, and communicate."
    ], size=16)


def add_usability_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(11), Inches(0.4),
                "Usability and workflow clarity", size=26, color=NAVY, bold=True)
    add_textbox(slide, Inches(0.7), Inches(0.92), Inches(10.8), Inches(0.55),
                "A major advantage of ProteoSphere is that the workflow is visible and explainable to collaborators who are not living inside the codebase.", size=17, color=SLATE)

    if SCREENSHOT_PATH.exists():
        slide.shapes.add_picture(str(SCREENSHOT_PATH), Inches(0.72), Inches(1.45), width=Inches(11.45))
    else:
        add_card(slide, Inches(0.72), Inches(1.45), Inches(11.45), Inches(4.8))

    add_card(slide, Inches(0.95), Inches(6.0), Inches(3.55), Inches(0.95))
    add_textbox(slide, Inches(1.15), Inches(6.22), Inches(3.1), Inches(0.18),
                "Guided workflow", size=15, color=TEAL, bold=True)
    add_textbox(slide, Inches(1.15), Inches(6.46), Inches(3.05), Inches(0.22),
                "Clear next-step guidance across data, models, inference, and outputs.", size=11, color=SLATE)

    add_card(slide, Inches(4.63), Inches(6.0), Inches(3.55), Inches(0.95))
    add_textbox(slide, Inches(4.83), Inches(6.22), Inches(3.1), Inches(0.18),
                "Platform literacy", size=15, color=BLUE, bold=True)
    add_textbox(slide, Inches(4.83), Inches(6.46), Inches(3.05), Inches(0.22),
                "The interface explains the pipeline rather than hiding it behind scripts.", size=11, color=SLATE)

    add_card(slide, Inches(8.31), Inches(6.0), Inches(3.55), Inches(0.95))
    add_textbox(slide, Inches(8.51), Inches(6.22), Inches(3.1), Inches(0.18),
                "Research communication", size=15, color=CORAL, bold=True)
    add_textbox(slide, Inches(8.51), Inches(6.46), Inches(3.05), Inches(0.22),
                "Makes it easier to walk collaborators through methods, assumptions, and outputs.", size=11, color=SLATE)


def add_comparison_slide(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.7), Inches(0.52), Inches(11), Inches(0.45),
                "Why this approach is stronger", size=26, color=NAVY, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.0), Inches(10.8), Inches(0.58),
                "ProteoSphere is designed to produce work that is broader, more reproducible, and easier to trust than the usual narrow-study pipeline.", size=18, color=SLATE)

    add_card(slide, Inches(0.9), Inches(1.95), Inches(4.9), Inches(4.65))
    add_textbox(slide, Inches(1.15), Inches(2.18), Inches(4.1), Inches(0.3),
                "Typical small-scale effort", size=18, color=CORAL, bold=True)
    add_bullets(slide, Inches(1.15), Inches(2.58), Inches(4.0), Inches(3.3), [
        "Small hand-built dataset",
        "Limited source coverage",
        "Weak provenance and conflict handling",
        "Random or weakly governed splits",
        "One graph or feature design treated as universal",
        "Difficult to present or reuse"
    ], size=16)

    add_card(slide, Inches(6.05), Inches(1.95), Inches(5.1), Inches(4.65), fill_color=MINT, line_color=TEAL)
    add_textbox(slide, Inches(6.3), Inches(2.18), Inches(4.3), Inches(0.3),
                "ProteoSphere approach", size=18, color=TEAL, bold=True)
    add_bullets(slide, Inches(6.3), Inches(2.58), Inches(4.2), Inches(3.3), [
        f"{metrics['tracked_sources']} staged source families with local-first packaging.",
        f"{metrics['bootstrap_records']:,} PDB records in the planning layer.",
        "Decision-grade field audit before policy use.",
        "Representative subset design and explicit split governance.",
        "Multiple graph designs and architecture paths.",
        "GUI-driven workflow built for research communication."
    ], size=16)


def add_close_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(5.95), prs.slide_width, Inches(1.55))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.2), Inches(0.6),
                "ProteoSphere", size=30, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.0), Inches(2.0), Inches(10.8), Inches(1.2),
                "A local-first platform for moving protein interaction ML from narrow, fragile studies toward scalable, reproducible, and explainable scientific workflows.", size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.2), Inches(3.5), Inches(10.4), Inches(0.7),
                "Core promise: broader source integration, stronger dataset governance, richer graph design, honest model selection, and a workflow that collaborators can actually follow.", size=18, color=SLATE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.9), Inches(6.35), Inches(11.4), Inches(0.4),
                "ProteoSphere | protein intelligence platform", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def build_deck() -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    metrics = load_metrics()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, metrics)
    add_problem_slide(prs)
    add_pipeline_slide(prs)
    add_graph_slide(prs)
    add_dataset_slide(prs, metrics)
    add_split_slide(prs, metrics)
    add_model_slide(prs)
    add_usability_slide(prs)
    add_comparison_slide(prs, metrics)
    add_close_slide(prs)

    output_path = OUTPUT_DIR / "ProteoSphere_platform_strategy_overview_2026-03-17.pptx"
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    print(build_deck())
